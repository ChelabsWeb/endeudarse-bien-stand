# -*- coding: utf-8 -*-
"""Genera 'Formación de formadores — Endeudarse bien (edición juego 2026)'.

Misma estructura y estética que el mazo original del programa Pin!
("Formación de formadores - Endeudarse bien.pptx", 68 slides), adaptado:
el bloque central del taller ahora es el juego "12 Meses" en equipos.

Uso:  python3 gen_formadores_pptx.py
Requiere: python-pptx. Opcional: /tmp/qr_juego.png y /tmp/qr_consulta.png.
Fuentes: Manual de tallerista Pin! 2025 (Taller 3) + el juego
https://chelabsweb.github.io/endeudarse-bien-stand/
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x25, 0x2A, 0x52)
VERDE  = RGBColor(0x3F, 0xB5, 0x77)
CREMA  = RGBColor(0xFA, 0xF3, 0xDC)
VIOLETA= RGBColor(0x6C, 0x63, 0xB5)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
ROJO   = RGBColor(0xE8, 0x5A, 0x66)
GRIS   = RGBColor(0x6B, 0x70, 0x93)
CELESTE= RGBColor(0x2E, 0x74, 0xB5)

DISP, BODY = 'Arial Black', 'Arial'
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
NPAG = [0]


def slide(bg=BLANCO):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    NPAG[0] += 1
    return s


def rbox(s, x, y, w, h, fill, line=None, lw=1.5):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    b.adjustments[0] = 0.08
    b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    b.shadow.inherit = False
    return b


def text(s, x, y, w, h, items, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = it.get('al', PP_ALIGN.LEFT)
        p.space_after = Pt(it.get('sa', 4)); p.space_before = Pt(it.get('sb', 0))
        for sg in it.get('segs', [it]):
            r = p.add_run(); r.text = sg['t']; f = r.font
            f.size = Pt(sg.get('s', it.get('s', 15)))
            f.bold = sg.get('b', it.get('b', False))
            f.name = sg.get('f', it.get('f', BODY))
            f.color.rgb = sg.get('c', it.get('c', NAVY))
    return tb


def dots(s, x, y, color=VERDE, n=5):
    for i in range(n):
        for j in range(2):
            shp = 'MATH_PLUS' if i < 2 else 'OVAL'
            d = s.shapes.add_shape(getattr(MSO_SHAPE, shp), Inches(x + j * 0.42), Inches(y + i * 0.42), Inches(0.16), Inches(0.16))
            d.fill.solid(); d.fill.fore_color.rgb = color; d.line.fill.background(); d.shadow.inherit = False


def pin(s, x, y, tc=CREMA):
    text(s, x, y, 2.5, 0.7, [
        {'t': 'Pin!', 's': 24, 'f': DISP, 'c': tc, 'sa': 0},
        {'t': 'Piques para manejar tu dinero', 's': 9, 'b': True, 'c': tc},
    ])


def head(s, kick, title, tc=NAVY, kc=VERDE):
    text(s, 0.6, 0.35, 12.2, 0.4, [{'t': kick, 's': 13, 'c': kc, 'b': True}])
    text(s, 0.55, 0.68, 12.3, 1.0, [{'t': title, 's': 31, 'c': tc, 'f': DISP}])


def footer(s, dark=False):
    text(s, 0.6, 7.1, 12.2, 0.3, [{'t': f'PIN! · TALLER 3 · ENDEUDARSE BIEN — EDICIÓN JUEGO 2026 · {NPAG[0]:02d}',
                                   's': 8.5, 'c': CREMA if dark else GRIS, 'b': True}])


def portada_seccion(kick, tit, sub=''):
    s = slide(NAVY)
    dots(s, 11.6, 0.6, VERDE); dots(s, 0.8, 5.2, VIOLETA)
    text(s, 1.0, 2.55, 11.3, 2.6, [
        {'t': kick, 's': 20, 'f': DISP, 'c': VERDE, 'sa': 8},
        {'t': tit, 's': 44, 'f': DISP, 'c': CREMA, 'sa': 8},
    ] + ([{'t': sub, 's': 16, 'b': True, 'c': BLANCO}] if sub else []))
    pin(s, 11.2, 6.3)
    footer(s, dark=True)
    return s


AGENDA = ['Actividad rompehielos', 'El juego: 12 meses', 'Huella financiera', 'Decidir bien', 'Ticket de salida']

def agenda(hl=-1):
    s = slide(NAVY)
    dots(s, 11.6, 0.6, VERDE)
    text(s, 0.9, 0.75, 8, 1.0, [{'t': 'AGENDA', 's': 40, 'f': DISP, 'c': CREMA}])
    y = 2.1
    for i, item in enumerate(AGENDA):
        act = (i == hl)
        ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y - 0.12), Inches(11.4), Pt(1.6))
        ln.fill.solid(); ln.fill.fore_color.rgb = VIOLETA; ln.line.fill.background(); ln.shadow.inherit = False
        if act:
            rbox(s, 0.8, y - 0.02, 11.7, 0.78, VERDE)
        text(s, 1.05, y + 0.06, 9.5, 0.6, [{'t': item, 's': 21, 'b': True, 'f': DISP if act else BODY,
                                            'c': NAVY if act else CREMA}])
        if act:
            text(s, 10.4, y + 0.1, 2.0, 0.5, [{'t': 'AHORA', 's': 14, 'f': DISP, 'c': NAVY, 'al': PP_ALIGN.RIGHT}])
        y += 0.95
    if hl != len(AGENDA) - 1:
        pin(s, 11.2, 6.3)
    footer(s, dark=True)
    return s


def barra(s, x, y, w_total, frac, label, val, color=VERDE, lw_label=4.4):
    text(s, x, y - 0.02, lw_label, 0.42, [{'t': label, 's': 12.5, 'b': True}])
    rbox(s, x + lw_label + 0.1, y, w_total, 0.32, CREMA, line=NAVY, lw=1.2)
    if frac > 0:
        f = rbox(s, x + lw_label + 0.1, y, max(w_total * frac, 0.15), 0.32, color, line=NAVY, lw=1.2)
    text(s, x + lw_label + 0.22 + w_total, y - 0.02, 0.9, 0.36, [{'t': str(val), 's': 15, 'b': True, 'f': DISP, 'c': color}])


QR1, QR2 = '/tmp/qr_juego.png', '/tmp/qr_consulta.png'

# ============================================================ 1 PORTADA
s = slide(NAVY)
dots(s, 11.6, 0.5, VERDE); dots(s, 0.8, 5.4, VIOLETA)
text(s, 0, 1.7, 13.333, 3.4, [
    {'t': 'Endeudarse bien', 's': 58, 'f': DISP, 'c': CREMA, 'al': PP_ALIGN.CENTER, 'sa': 6},
    {'t': 'decisiones financieras informadas', 's': 22, 'b': True, 'c': VERDE, 'al': PP_ALIGN.CENTER, 'sa': 10},
    {'t': 'Taller 3 · Formación de formadores · edición juego 2026', 's': 15, 'b': True, 'c': BLANCO, 'al': PP_ALIGN.CENTER},
])
rbox(s, 3.87, 5.25, 5.6, 0.65, VERDE)
text(s, 3.87, 5.35, 5.6, 0.5, [{'t': 'EL TALLER AHORA SE JUEGA', 's': 16, 'f': DISP, 'c': NAVY, 'al': PP_ALIGN.CENTER}])
text(s, 0, 6.15, 13.333, 0.4, [{'t': 'Pin! · INJU · BCU · CAF — Central para Vos', 's': 12, 'b': True, 'c': CREMA, 'al': PP_ALIGN.CENTER}])
footer(s, dark=True)

# ============================================================ 2 AGENDA
agenda()

# ============================================================ ROMPEHIELO
agenda(0)

s = slide(); head(s, 'DESAFÍO FINANCIERO', 'El desafío')
text(s, 0.6, 1.75, 12.1, 1.5, [
    {'t': 'Objetivo', 's': 15, 'b': True, 'c': CELESTE, 'sa': 3},
    {'t': 'Fomentar la reflexión sobre decisiones financieras mediante un juego dinámico y participativo, con datos reales sobre cómo se endeudan las familias uruguayas y sus consecuencias.', 's': 14, 'sa': 8},
])
b = rbox(s, 0.6, 3.15, 5.9, 3.2, CREMA, line=VERDE, lw=2)
text(s, 0.9, 3.35, 5.3, 2.8, [
    {'t': 'Materiales', 's': 15, 'b': True, 'c': CELESTE, 'sa': 6},
    {'t': '• Proyector con esta presentación (y luego el tablero del juego).', 's': 13, 'sa': 5},
    {'t': '• Pizarra para anotar los puntos de cada equipo.', 's': 13, 'sa': 5},
    {'t': '• Cronómetro o reloj visible.', 's': 13, 'sa': 5},
    {'t': '• Las 3 máquinas del juego, ya conectadas a la sala.', 's': 13, 'b': True, 'sa': 5},
])
b = rbox(s, 6.85, 3.15, 5.9, 3.2, NAVY)
text(s, 7.15, 3.35, 5.3, 2.8, [
    {'t': 'Duración', 's': 15, 'b': True, 'c': VERDE, 'sa': 6},
    {'t': '15 minutos', 's': 26, 'f': DISP, 'c': CREMA, 'sa': 8},
    {'t': 'Este rompehielo además arma los equipos que van a jugar "12 Meses".', 's': 13.5, 'b': True, 'c': BLANCO},
])
footer(s)

s = slide(); head(s, 'DESAFÍO FINANCIERO', 'Instrucciones')
text(s, 0.6, 1.8, 12.1, 4.6, [
    {'t': 'Se divide el grupo en 3 equipos (los mismos que después van a jugar, ~6-7 por equipo).', 's': 16, 'b': True, 'sa': 10},
    {'t': 'Al estilo de "100 uruguayos dicen": si la respuesta coincide con una de la lista oficial, suman los puntos de esa respuesta. Los puntos reflejan cuántos uruguayos están en cada situación.', 's': 15, 'sa': 10},
    {'t': 'Los equipos se alternan. Cada equipo tiene 30 segundos para discutir y solo su vocero/a da la respuesta final.', 's': 15, 'sa': 10},
    {'t': 'Gana el equipo que suma más puntos… y se lleva la primera elección de máquina y nombre de equipo.', 's': 15, 'b': True, 'c': VERDE, 'sa': 10},
])
footer(s)

s = slide(); head(s, 'DESAFÍO FINANCIERO', '¿Quién empieza?')
b = rbox(s, 0.6, 1.9, 12.1, 3.0, CREMA, line=VERDE, lw=2)
text(s, 1.0, 2.25, 11.3, 2.4, [
    {'t': 'En Uruguay hay algo más de 600.000 jóvenes de entre 18 y 30 años.', 's': 19, 'b': True, 'sa': 8},
    {'t': '¿Cuántos creen ustedes que han solicitado un préstamo a un banco o financiera?', 's': 22, 'f': DISP, 'c': CELESTE},
])
text(s, 0.6, 5.3, 12.1, 1.0, [
    {'t': 'Cada equipo tiene 30 segundos para acordar un número. El que más se acerque, empieza.', 's': 15, 'b': True, 'c': GRIS},
])
footer(s)

s = slide(NAVY); dots(s, 11.6, 0.6, VERDE)
text(s, 0.9, 0.7, 11, 0.9, [{'t': 'Respuesta correcta', 's': 24, 'f': DISP, 'c': VERDE}])
text(s, 0, 2.1, 13.333, 2.2, [
    {'t': '297.329', 's': 100, 'f': DISP, 'c': CREMA, 'al': PP_ALIGN.CENTER, 'sa': 4},
    {'t': '1 DE CADA 2 JÓVENES YA PASÓ POR VENTANILLA', 's': 20, 'f': DISP, 'c': VERDE, 'al': PP_ALIGN.CENTER},
])
text(s, 1.5, 5.3, 10.3, 1.3, [
    {'t': 'Registrados en la Central de Riesgos del BCU. Ojo: estar registrado no significa no estar pagando.', 's': 14.5, 'b': True, 'c': BLANCO, 'al': PP_ALIGN.CENTER, 'sa': 5},
    {'t': '"Si fuéramos como el promedio, un equipo y medio de esta sala estaría endeudado." ¿Les sorprende?', 's': 14.5, 'b': True, 'c': VERDE, 'al': PP_ALIGN.CENTER},
])
footer(s, dark=True)

s = slide(); head(s, 'PREGUNTA 1', '¿Cuáles son las principales razones de endeudamiento en los hogares uruguayos?')
text(s, 0.6, 1.95, 12.1, 0.8, [{'t': 'Las razones pueden ir desde cubrir gastos básicos hasta inversiones grandes. Su tarea: adivinar las más comunes. Solo suman si está en la lista.', 's': 14}])
razones = ['arreglar la casa', 'cancelar una deuda', 'comprar un vehículo o muebles', 'comprar una vivienda',
           'gastos corrientes (comida, ropa, limpieza)', 'gastos mensuales (UTE, OSE, alquiler)', 'inversión en un negocio', 'irse de vacaciones']
for i, rz in enumerate(razones):
    x = 0.6 + (i % 2) * 6.25; y = 2.95 + (i // 2) * 0.85
    rbox(s, x, y, 5.95, 0.68, CREMA, line=VIOLETA, lw=1.5)
    text(s, x + 0.25, y + 0.12, 5.5, 0.5, [{'t': rz, 's': 14, 'b': True}])
footer(s)

s = slide(); head(s, 'PREGUNTA 1 · TABLERO', 'Las respuestas de los uruguayos')
datos = [('Gastos corrientes (comida, ropa, hogar)', 25, VERDE), ('Gastos mensuales (UTE, OSE, alquiler)', 18, CELESTE),
         ('Cancelar una deuda previa', 17, VIOLETA), ('Vehículo o muebles', 15, NAVY),
         ('Arreglo de la vivienda', 12, GRIS), ('Compra de vivienda', 5, GRIS),
         ('Inversión en un negocio', 3, GRIS), ('Vacaciones', 2, GRIS)]
y = 1.95
for label, val, col in datos:
    barra(s, 0.7, y, 6.4, val / 25.0, label, val, col)
    y += 0.5
text(s, 0.6, 6.15, 12.1, 0.8, [{'segs': [
    {'t': 'Cuando la deuda es para el día a día, lo que falta no es crédito: ', 's': 14.5, 'b': True},
    {'t': 'son más ingresos o menos gastos.', 's': 14.5, 'b': True, 'c': ROJO},
]}])
footer(s)

s = slide(); head(s, 'PREGUNTA 2', '¿Qué tan buenos deudores son los jóvenes uruguayos?')
text(s, 0.6, 1.9, 12.1, 0.7, [{'t': 'La Central de Riesgos del BCU clasifica a cada deudor en 5 niveles según los días de atraso:', 's': 14.5, 'b': True}])
cats = [('1', 'Capacidad de pago fuerte: al día o atraso menor a 10 días', VERDE),
        ('2', 'Capacidad adecuada o problemas potenciales: 10 a 60 días', CELESTE),
        ('3', 'Capacidad comprometida: 60 a 180 días', VIOLETA),
        ('4', 'Varias deudas, al menos una irrecuperable (más de 180 días)', ROJO),
        ('5', 'TODAS sus deudas irrecuperables', NAVY)]
y = 2.65
for n, d, col in cats:
    c = rbox(s, 0.7, y, 0.6, 0.6, col)
    text(s, 0.7, y + 0.08, 0.6, 0.45, [{'t': n, 's': 20, 'f': DISP, 'c': BLANCO, 'al': PP_ALIGN.CENTER}])
    text(s, 1.5, y + 0.08, 11.2, 0.5, [{'t': d, 's': 14.5, 'b': True}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.78
text(s, 0.6, 6.6, 12.1, 0.5, [{'t': '¿Cómo creen que están repartidos los jóvenes? Los equipos se alternan nombrando categorías.', 's': 13.5, 'b': True, 'c': GRIS}])
footer(s)

s = slide(); head(s, 'PREGUNTA 2 · TABLERO', 'Así están repartidos (de cada 100 jóvenes deudores)')
datos2 = [('CAT 1 · Pago fuerte', 50, VERDE), ('CAT 5 · Todo irrecuperable', 30, ROJO),
          ('CAT 4 · Al menos una irrecuperable', 10, VIOLETA), ('CAT 2 · Pago adecuado', 7, CELESTE),
          ('CAT 3 · Pago comprometido', 3, GRIS)]
y = 2.0
for label, val, col in datos2:
    barra(s, 0.7, y, 6.4, val / 50.0, label, val, col)
    y += 0.56
b = rbox(s, 0.6, 5.15, 12.1, 1.35, ROJO)
text(s, 0.95, 5.32, 11.5, 1.05, [
    {'t': '1 DE CADA 3 YA TIENE TODAS SUS DEUDAS INCOBRABLES', 's': 20, 'f': DISP, 'c': BLANCO, 'sa': 4},
    {'t': 'La mitad paga impecable — y un tercio arrancó con la huella quemada 5 años. Esa diferencia se juega en decisiones chicas.', 's': 13.5, 'b': True, 'c': BLANCO},
])
footer(s)

s = slide(); head(s, 'DESAFÍO FINANCIERO', 'Disparadores para la reflexión')
b = rbox(s, 0.6, 1.8, 5.95, 4.5, CREMA, line=VERDE, lw=2)
text(s, 0.9, 2.0, 5.4, 4.1, [
    {'t': 'Sobre los motivos (P1)', 's': 15, 'b': True, 'c': CELESTE, 'sa': 7},
    {'t': '• ¿Qué motivo les sorprendió más? ¿Por qué?', 's': 13.5, 'sa': 6},
    {'t': '• ¿Por qué "vacaciones" o "invertir" están tan abajo?', 's': 13.5, 'sa': 6},
    {'t': '• ¿Se puede evitar endeudarse por comida o servicios? ¿Hay que hacerlo?', 's': 13.5, 'sa': 6},
])
b = rbox(s, 6.8, 1.8, 5.95, 4.5, CREMA, line=VIOLETA, lw=2)
text(s, 7.1, 2.0, 5.4, 4.1, [
    {'t': 'Sobre las categorías (P2)', 's': 15, 'b': True, 'c': CELESTE, 'sa': 7},
    {'t': '• ¿Qué piensan de que la mayoría esté al día?', 's': 13.5, 'sa': 6},
    {'t': '• ¿Y de que 1 de cada 3 menores de 30 tenga todo irrecuperable?', 's': 13.5, 'sa': 6},
    {'t': '• ¿Cómo afectan estas clasificaciones el acceso a crédito futuro? (puente a "huella financiera")', 's': 13.5, 'b': True, 'sa': 6},
])
footer(s)

s = portada_seccion('MENSAJE CLAVE', 'El endeudamiento no siempre\nes malo…',
                    'pero puede reflejar falta de opciones o de planificación. Un mal endeudamiento afecta la huella financiera y limita las opciones futuras.')

# ============================================================ EL JUEGO
agenda(1)

s = slide(); head(s, 'EL JUEGO', 'Ahora ustedes son Joaquín')
text(s, 0.6, 1.75, 12.1, 1.1, [
    {'t': 'Joaquín gana $32.000, debe media vida y quiere una moto para hacer delivery.', 's': 18, 'b': True, 'sa': 5},
    {'t': 'Van a manejarle la plata un año entero: 12 meses, una decisión por mes.', 's': 16, 'b': True, 'c': CELESTE},
])
hud = [('BILLETERA', 'la plata del mes', VERDE), ('FONDO', 'el colchón de emergencia', CELESTE),
       ('INTERESES REGALADOS', 'lo pagado de más por financiar', ROJO),
       ('RACHA', 'decisiones buenas seguidas = bonus', VIOLETA),
       ('HUELLA', 'cada marca queda 5 años', NAVY)]
y = 3.05
for t, d, col in hud:
    rbox(s, 0.6, y, 3.6, 0.55, col)
    text(s, 0.6, y + 0.1, 3.6, 0.4, [{'t': t, 's': 12.5, 'f': DISP, 'c': BLANCO, 'al': PP_ALIGN.CENTER}])
    text(s, 4.45, y + 0.08, 8.2, 0.45, [{'t': d, 's': 14, 'b': True}], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.68
text(s, 0.6, 6.5, 12.1, 0.5, [{'t': 'En el camino: quizzes con los datos que acaban de ver, una bocha inesperada en julio y el aguinaldo en diciembre.', 's': 13.5, 'b': True, 'c': GRIS}])
footer(s)

s = slide(); head(s, 'EL JUEGO', 'Cómo jugamos')
pasos = [('1', 'EQUIPOS', 'Cada equipo, a su máquina. Elijan nombre (hasta 12 letras): así aparecen en el tablero.', VERDE),
         ('2', 'CÓDIGO DE SALA', 'Las 3 máquinas ponen el MISMO código de 4 letras. Mismos eventos para todos: la comparación es justa.', CELESTE),
         ('3', 'EL TABLERO', 'El proyector muestra la carrera en vivo: puntos, mes por mes, y la huella de cada equipo.', VIOLETA)]
y = 1.85
for n, t, d, col in pasos:
    b = rbox(s, 0.6, y, 12.1, 1.35, CREMA, line=col, lw=2)
    c = rbox(s, 0.85, y + 0.3, 0.75, 0.75, col)
    text(s, 0.85, y + 0.4, 0.75, 0.6, [{'t': n, 's': 24, 'f': DISP, 'c': BLANCO, 'al': PP_ALIGN.CENTER}])
    text(s, 1.9, y + 0.18, 10.6, 1.05, [
        {'t': t, 's': 15, 'f': DISP, 'sa': 3},
        {'t': d, 's': 13.5, 'b': True},
    ])
    y += 1.55
text(s, 0.6, 6.5, 12.1, 0.5, [{'t': 'Reglas de oro: 45 segundos por decisión (hay reloj) · se debate en equipo · decide el vocero.', 's': 14, 'b': True, 'c': ROJO}])
footer(s)

s = slide(); head(s, 'EL JUEGO', 'Los roles del equipo')
roles = [('VOCERO/A', 'El único que toca la pantalla. Rota cada 3 meses del juego.', VERDE),
         ('CONTADOR/A', 'Hace la cuenta EN VOZ ALTA antes de decidir: cuota × cuotas.', CELESTE),
         ('ABOGADO/A DEL DIABLO', 'Defiende la opción tentadora. Si nadie tienta, no hay mérito en resistir.', VIOLETA)]
y = 1.9
for t, d, col in roles:
    b = rbox(s, 0.6, y, 12.1, 1.3, col)
    text(s, 1.0, y + 0.18, 11.3, 1.0, [
        {'t': t, 's': 17, 'f': DISP, 'c': BLANCO, 'sa': 4},
        {'t': d, 's': 14.5, 'b': True, 'c': BLANCO},
    ])
    y += 1.5
text(s, 0.6, 6.45, 12.1, 0.6, [{'t': 'Los demás: opinan, discuten y gritan lo justo. Todos deciden, uno toca.', 's': 15, 'b': True, 'c': GRIS}])
footer(s)

s = slide(NAVY); dots(s, 11.6, 0.6, VERDE); dots(s, 0.8, 5.2, VIOLETA)
text(s, 0, 1.5, 13.333, 1.4, [{'t': '¡A JUGAR!', 's': 64, 'f': DISP, 'c': VERDE, 'al': PP_ALIGN.CENTER}])
text(s, 0, 3.1, 13.333, 0.6, [{'t': 'chelabsweb.github.io/endeudarse-bien-stand', 's': 22, 'b': True, 'c': CREMA, 'al': PP_ALIGN.CENTER}])
rbox(s, 4.42, 4.0, 4.5, 1.1, CREMA)
text(s, 4.42, 4.12, 4.5, 0.9, [
    {'t': 'CÓDIGO DE SALA', 's': 13, 'b': True, 'c': GRIS, 'al': PP_ALIGN.CENTER, 'sa': 2},
    {'t': '_ _ _ _', 's': 30, 'f': DISP, 'c': NAVY, 'al': PP_ALIGN.CENTER},
])
text(s, 0, 5.4, 13.333, 0.9, [
    {'t': '15 minutos · 12 meses · el tablero queda en el proyector', 's': 16, 'b': True, 'c': BLANCO, 'al': PP_ALIGN.CENTER, 'sa': 4},
    {'t': '(el facilitador cambia esta pantalla por el modo Proyector del juego)', 's': 12.5, 'b': True, 'c': VERDE, 'al': PP_ALIGN.CENTER},
])
footer(s, dark=True)

s = slide(); head(s, 'DESPUÉS DE LA PARTIDA', 'Leemos el tablero juntos')
b = rbox(s, 0.6, 1.8, 5.95, 4.5, CREMA, line=VERDE, lw=2)
text(s, 0.9, 2.0, 5.4, 4.1, [
    {'t': 'EL PODIO', 's': 15, 'b': True, 'c': CELESTE, 'sa': 7},
    {'t': '• Intereses regalados = plata quemada: ¿cuánto regaló cada equipo?', 's': 13.5, 'sa': 6},
    {'t': '• Marcas en la huella = 5 años: ¿cuáles se podían evitar?', 's': 13.5, 'sa': 6},
    {'t': '• Categoría BCU final: conectar con el tablero de la pregunta 2.', 's': 13.5, 'sa': 6},
])
b = rbox(s, 6.8, 1.8, 5.95, 4.5, CREMA, line=ROJO, lw=2)
text(s, 7.1, 2.0, 5.4, 4.1, [
    {'t': '"DONDE MÁS NOS CLAVAMOS"', 's': 15, 'b': True, 'c': ROJO, 'sa': 7},
    {'t': 'El proyector muestra los errores más repetidos entre los equipos, con su explicación ("la posta").', 's': 13.5, 'sa': 6},
    {'t': 'Por cada error repetido: ¿por qué era tentador? ¿qué señal lo delataba?', 's': 13.5, 'b': True, 'sa': 6},
    {'t': 'Cada equipo cuenta su peor decisión y qué haría distinto.', 's': 13.5, 'b': True, 'sa': 6},
])
footer(s)

s = slide(); head(s, 'LO QUE VIVIERON EN EL JUEGO', 'Las apariencias engañan')
text(s, 0.6, 1.7, 12.1, 0.6, [{'t': 'El celular de $4.500 (como el de Juana en el quiz): tres formas de pagarlo.', 's': 15, 'b': True}])
cols = [('CONTADO', '1 × $4.500', '$4.500', VERDE, 'Puede pedir descuento y pagar menos'),
        ('CUOTAS', '3 × $1.500', '$4.500', CELESTE, 'Mismo total… ¿entonces quién paga la espera?'),
        ('CUOTAS CON INTERESES', '12 × $500', '$6.000', ROJO, 'Paga $1.500 más: un tercio de celular extra')]
for i, (t, cu, tot, col, nota) in enumerate(cols):
    x = 0.6 + i * 4.18
    b = rbox(s, x, 2.5, 3.9, 3.3, CREMA, line=col, lw=2.5)
    text(s, x + 0.2, 2.7, 3.5, 3.0, [
        {'t': t, 's': 14, 'f': DISP, 'c': col, 'sa': 6, 'al': PP_ALIGN.CENTER},
        {'t': cu, 's': 17, 'b': True, 'al': PP_ALIGN.CENTER, 'sa': 6},
        {'t': tot, 's': 30, 'f': DISP, 'c': col, 'al': PP_ALIGN.CENTER, 'sa': 8},
        {'t': nota, 's': 12.5, 'b': True, 'al': PP_ALIGN.CENTER},
    ])
text(s, 0.6, 6.15, 12.1, 0.7, [{'t': 'Si alguien espera 3 meses para cobrar $4.500… otro debería estar dispuesto a dártelo por menos hoy. Y si Juana puede pagar $1.500 por mes: ¿por qué no ahorrarlos ANTES?', 's': 13.5, 'b': True, 'c': GRIS}])
footer(s)

s = slide(); head(s, 'LO QUE VIVIERON EN EL JUEGO', 'Lo barato sale caro: el pago mínimo')
b = rbox(s, 0.6, 1.9, 5.95, 3.9, CREMA, line=VERDE, lw=2.5)
text(s, 0.9, 2.15, 5.4, 3.4, [
    {'t': 'PAGO TOTAL (o plan de ataque)', 's': 15, 'f': DISP, 'c': VERDE, 'sa': 8},
    {'t': 'La deuda SE ACHICA', 's': 22, 'f': DISP, 'sa': 8},
    {'t': 'Unos meses de cinturón apretado y la tarjeta muere. Intereses mínimos.', 's': 13.5, 'b': True},
])
b = rbox(s, 6.8, 1.9, 5.95, 3.9, CREMA, line=ROJO, lw=2.5)
text(s, 7.1, 2.15, 5.4, 3.4, [
    {'t': 'PAGO MÍNIMO', 's': 15, 'f': DISP, 'c': ROJO, 'sa': 8},
    {'t': 'La deuda SE AGRANDA', 's': 22, 'f': DISP, 'sa': 8},
    {'t': 'El mínimo apenas cubre intereses: pagás y pagás y el saldo ni se entera. La deuda eterna.', 's': 13.5, 'b': True},
])
text(s, 0.6, 6.1, 12.1, 0.6, [{'t': 'En el juego lo vivieron en el Mes 2. En la vida real, el resumen de la tarjeta lo muestra todos los meses.', 's': 14, 'b': True, 'c': GRIS}])
footer(s)

s = slide(); head(s, 'EL MÉTODO', 'PAMY: cuatro preguntas antes de firmar')
pamy = [('P', 'PRECISO', '¿Lo preciso de verdad? ¿Y lo preciso YA? Un préstamo es una decisión de tiempo: ¿qué pasa si espero?', VERDE),
        ('A', 'AHORRO', '¿Tengo ahorrado? ¿Voy a poder ahorrar? Si puedo pagar la cuota, puedo ser mi propio prestamista.', CELESTE),
        ('M', 'MEJORES', '¿Cuál es mi mejor opción? Bancos (ojo BROU), financieras, cooperativas, tarjetas. Comparar el TOTAL, no la cuota.', VIOLETA),
        ('Y', 'Y DESPUÉS', '¿Cómo lo voy a pagar? ¿Qué gasto comprimo? ¿Cuál es el plan para sostenerlo?', NAVY)]
for i, (l, t, d, col) in enumerate(pamy):
    x = 0.6 + (i % 2) * 6.25; y = 1.85 + (i // 2) * 2.3
    b = rbox(s, x, y, 5.95, 2.05, col)
    text(s, x + 0.3, y + 0.25, 1.1, 1.4, [{'t': l, 's': 44, 'f': DISP, 'c': BLANCO}])
    text(s, x + 1.5, y + 0.28, 4.3, 1.6, [
        {'t': t, 's': 16, 'f': DISP, 'c': BLANCO, 'sa': 4},
        {'t': d, 's': 12, 'b': True, 'c': BLANCO},
    ])
footer(s)

s = slide(); head(s, 'LA CUENTA QUE SALVA', '¿Cuántas cuotas son para vos?')
ej = [('$15.000 en 6 cuotas de $2.999', '6 × $2.999 = $17.994', '5 cuotas para vos · 1 entera para la financiera', VERDE),
      ('$80.000 en 36 cuotas de $3.999', '36 × $3.999 = $143.964', '20 cuotas para vos · 16 enteras para la financiera', ROJO)]
y = 1.9
for t, cta, rep, col in ej:
    b = rbox(s, 0.6, y, 12.1, 1.8, CREMA, line=col, lw=2.5)
    text(s, 1.0, y + 0.18, 11.3, 1.5, [
        {'t': t, 's': 15, 'b': True, 'sa': 3},
        {'t': cta, 's': 24, 'f': DISP, 'c': col, 'sa': 3},
        {'t': rep, 's': 14.5, 'b': True},
    ])
    y += 2.05
text(s, 0.6, 6.1, 12.1, 0.8, [{'t': 'La TEA (tasa efectiva anual) dice cuánto interés pagás por año — se publica en 5 lugares. Pero la cuenta que nunca falla es cuota × cuotas, pedida POR ESCRITO.', 's': 14, 'b': True, 'c': GRIS}])
footer(s)

s = portada_seccion('MENSAJE CLAVE', 'Endeudarse es ahorrar…',
                    'pero ahorrar no solamente para vos: también para pagarle los intereses a quien te prestó. Evaluar costo total, intereses y plazos es la base de una decisión responsable.')

# ============================================================ HUELLA
agenda(2)

s = slide(); head(s, 'HUELLA FINANCIERA', 'Los dos espejos donde te miran')
b = rbox(s, 0.6, 1.8, 5.95, 4.3, CELESTE)
text(s, 0.9, 2.0, 5.4, 3.9, [
    {'t': 'BCU · CENTRAL DE RIESGOS', 's': 15, 'f': DISP, 'c': BLANCO, 'sa': 7},
    {'t': '• Solo sistema financiero regulado: bancos, financieras, cooperativas.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
    {'t': '• Tu nota del 1 al 5 según los días de atraso.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
    {'t': '• No borra la información aunque canceles la deuda.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
])
b = rbox(s, 6.8, 1.8, 5.95, 4.3, VIOLETA)
text(s, 7.1, 2.0, 5.4, 3.9, [
    {'t': 'CLEARING DE INFORMES', 's': 15, 'f': DISP, 'c': BLANCO, 'sa': 7},
    {'t': '• Servicios (UTE, OSE), comercios, cheques rechazados, refinanciaciones.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
    {'t': '• Muestra quién te consultó en los últimos 6 meses.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
    {'t': '• Las marcas duran 5 años desde el registro.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
])
text(s, 0.6, 6.3, 12.1, 0.6, [{'t': 'EL BANCO MIRA LOS DOS. En el juego: por eso la UTE impaga también marcó la huella.', 's': 16, 'f': DISP, 'al': PP_ALIGN.CENTER}])
footer(s)

s = slide(); head(s, 'HUELLA FINANCIERA', 'Consultá la tuya: es gratis')
b = rbox(s, 0.6, 1.85, 5.95, 4.3, CREMA, line=CELESTE, lw=2)
if os.path.exists(QR2):
    s.shapes.add_picture(QR2, Inches(1.0), Inches(2.2), Inches(2.1), Inches(2.1))
text(s, 3.3, 2.4, 3.1, 2.0, [
    {'t': 'BCU', 's': 16, 'f': DISP, 'c': CELESTE, 'sa': 5},
    {'t': 'consultadeuda.bcu.gub.uy', 's': 14, 'b': True, 'sa': 5},
    {'t': 'Tu categoría y tus deudas en el sistema financiero.', 's': 12, 'b': True, 'c': GRIS},
])
b = rbox(s, 6.8, 1.85, 5.95, 4.3, CREMA, line=VIOLETA, lw=2)
text(s, 7.2, 2.4, 5.2, 2.6, [
    {'t': 'CLEARING', 's': 16, 'f': DISP, 'c': VIOLETA, 'sa': 5},
    {'t': 'clearing.com.uy/personas', 's': 14, 'b': True, 'sa': 5},
    {'t': 'Servicios, cheques, refinanciaciones y quién te consultó. La primera consulta del período no tiene costo.', 's': 12, 'b': True, 'c': GRIS},
])
footer(s)

s = slide(); head(s, 'CASO', 'La huella de Joaquín')
text(s, 0.6, 1.75, 12.1, 0.7, [{'t': 'Joaquín, 27 años, técnico de mantenimiento en Canelones. Gana $32.000 por mes.', 's': 15.5, 'b': True}])
b = rbox(s, 0.6, 2.5, 5.95, 3.9, CREMA, line=CELESTE, lw=2)
text(s, 0.9, 2.7, 5.4, 3.5, [
    {'t': 'Historial financiero', 's': 14, 'b': True, 'c': CELESTE, 'sa': 6},
    {'t': '• Tarjeta con saldo de $10.000: paga el mínimo hace 8 meses (sin atraso).', 's': 12.5, 'sa': 5},
    {'t': '• Préstamo de $120.000 por el auto hace 3 años: le quedan $20.000, al día.', 's': 12.5, 'sa': 5},
    {'t': '• UTE atrasada 2 meses ($4.000).', 's': 12.5, 'sa': 5},
    {'t': '• Un cheque rechazado que sustituyó.', 's': 12.5, 'sa': 5},
    {'t': '• Sin deudas irrecuperables.', 's': 12.5, 'sa': 5},
])
b = rbox(s, 6.8, 2.5, 5.95, 3.9, CREMA, line=VIOLETA, lw=2)
text(s, 7.1, 2.7, 5.4, 3.5, [
    {'t': 'Acciones recientes', 's': 14, 'b': True, 'c': VIOLETA, 'sa': 6},
    {'t': '• Refinanció la deuda de la tarjeta con una financiera.', 's': 12.5, 'sa': 5},
    {'t': '• Canceló un préstamo menor de $5.000 hace un año.', 's': 12.5, 'sa': 5},
    {'t': '• 6 consultas a su Clearing en 6 meses (Anda, Itaú, Motociclo, Creditel ×2, BROU).', 's': 12.5, 'sa': 10},
    {'t': 'Quiere pedir $50.000 en el banco para la moto del delivery.', 's': 13.5, 'b': True, 'c': NAVY},
])
footer(s)

s = slide(); head(s, 'CASO', '¿Le prestarían?')
b = rbox(s, 0.6, 1.8, 5.95, 4.3, CELESTE)
text(s, 0.9, 2.0, 5.4, 3.9, [
    {'t': 'EL BCU DICE…', 's': 15, 'f': DISP, 'c': BLANCO, 'sa': 7},
    {'t': 'Tarjeta al día (paga el mínimo, pero sin atraso).', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
    {'t': 'Préstamo del auto al día.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
    {'t': 'La UTE y el cheque NO entran acá.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 8},
    {'t': '→ Categoría 1: parece que SÍ.', 's': 17, 'f': DISP, 'c': CREMA},
])
b = rbox(s, 6.8, 1.8, 5.95, 4.3, VIOLETA)
text(s, 7.1, 2.0, 5.4, 3.9, [
    {'t': 'EL CLEARING DICE…', 's': 15, 'f': DISP, 'c': BLANCO, 'sa': 7},
    {'t': 'Operación incumplida: SÍ (UTE, 60 días).', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
    {'t': 'Refinanciación: SÍ (la tarjeta).', 's': 13, 'b': True, 'c': BLANCO, 'sa': 5},
    {'t': 'Cheque rechazado: SÍ. Consultas: 6 en 6 meses.', 's': 13, 'b': True, 'c': BLANCO, 'sa': 8},
    {'t': '→ Mmm… acá la cosa cambia.', 's': 17, 'f': DISP, 'c': CREMA},
])
text(s, 0.6, 6.3, 12.1, 0.6, [{'t': 'Debate y votación: ¿le prestan? ¿Pesa más la nota del BCU o el Clearing? ¿Importa que la moto genere ingresos?', 's': 14, 'b': True, 'c': GRIS, 'al': PP_ALIGN.CENTER}])
footer(s)

s = portada_seccion('MENSAJE CLAVE', 'La huella es tu reputación',
                    'Cuidarla da acceso a mejores oportunidades y condiciones. Cada decisión cuenta: pagar a tiempo, conocer el historial y planificar con responsabilidad.')

# ============================================================ DECIDIR BIEN
agenda(3)

s = slide(); head(s, 'DECIDIR BIEN', 'La decisión de Joaquín, en números')
cols = [('HOY', ['Sueldo: $32.000', 'Gastos: −$24.700', 'Deudas: −$7.000', 'Neto: $300', 'Sin fondo de emergencia', 'Debe $4.000 de UTE'], GRIS),
        ('EN 4 MESES', ['Sueldo: $32.000', 'Gastos: −$24.700', 'Deudas: −$2.000', 'Neto: $4.300', 'Sin fondo de emergencia', 'Debe $4.000 de UTE'], CELESTE),
        ('CON PRÉSTAMO', ['Sueldo: $32.000', 'PedidosYa: +$7.000', 'Gastos: −$24.700 · Deudas: −$5.999', 'Neto: $7.301', 'Fondo de emergencia: $6.000', 'Sin deuda de UTE'], VERDE)]
for i, (t, filas, col) in enumerate(cols):
    x = 0.6 + i * 4.18
    b = rbox(s, x, 1.85, 3.9, 4.35, CREMA, line=col, lw=2.5)
    items = [{'t': t, 's': 16, 'f': DISP, 'c': col, 'al': PP_ALIGN.CENTER, 'sa': 8}]
    for f in filas:
        items.append({'t': f, 's': 12.8, 'b': True, 'al': PP_ALIGN.CENTER, 'sa': 5})
    text(s, x + 0.15, 2.05, 3.6, 4.0, items)
text(s, 0.6, 6.4, 12.1, 0.5, [{'t': 'Préstamo: $60.000 · cuota $3.999 · la moto suma ~$7.000/mes de delivery… si todo sale bien.', 's': 13.5, 'b': True, 'c': GRIS, 'al': PP_ALIGN.CENTER}])
footer(s)

s = slide(); head(s, 'DECIDIR BIEN', 'Ahora ustedes son Joaquín (otra vez)')
text(s, 0.6, 1.8, 12.1, 1.5, [
    {'t': 'La consigna: contarle el plan a la IA asesora y dejarse hacer preguntas.', 's': 17, 'b': True, 'sa': 8},
    {'t': 'Antes de empezar, cada uno anota:', 's': 15, 'b': True, 'c': CELESTE},
])
b = rbox(s, 0.6, 3.3, 12.1, 1.6, CREMA, line=VERDE, lw=2)
text(s, 1.0, 3.5, 11.3, 1.2, [
    {'t': '¿Tomarían el préstamo?  SÍ / NO', 's': 18, 'f': DISP, 'sa': 5},
    {'t': 'Y del 1 al 7: ¿qué tan convencidos están de que es una buena decisión?', 's': 15, 'b': True},
])
text(s, 0.6, 5.2, 12.1, 1.2, [
    {'t': 'La IA no les va a decir qué hacer: pregunta cómo van a cubrir las cuotas, qué pasa si surge un imprevisto, y cómo se imaginan la vida en un año con esa decisión.', 's': 14, 'b': True, 'c': GRIS},
])
footer(s)

s = slide(); head(s, 'DECIDIR BIEN', 'Después de hablar con la IA')
b = rbox(s, 0.6, 1.9, 12.1, 3.6, CREMA, line=VIOLETA, lw=2)
text(s, 1.0, 2.2, 11.3, 3.0, [
    {'t': '¿Mantienen la decisión?', 's': 22, 'f': DISP, 'sa': 8},
    {'t': 'Indiquen de nuevo, del 1 al 7, qué tan convencidos están.', 's': 16, 'b': True, 'sa': 8},
    {'t': '¿Cambió algo? ¿Qué apareció en la conversación que no habían pensado?', 's': 16, 'b': True, 'c': VIOLETA},
])
text(s, 0.6, 5.8, 12.1, 0.8, [{'t': 'El punto no es la respuesta: es descubrir que reflexionar ANTES de firmar cambia la decisión.', 's': 15, 'b': True, 'c': GRIS}])
footer(s)

s = portada_seccion('MENSAJE CLAVE', 'Los números no alcanzan',
                    'Tomar decisiones financieras requiere introspección y planificación. Reflexionar antes de actuar marca la diferencia entre una decisión que te empodera y una que te limita.')

# ============================================================ CIERRE
agenda(4)

s = slide(); head(s, 'MINUTO DE ORO', 'Las ideas clave del taller')
ideas = [('EL DESAFÍO', 'La mitad de los jóvenes ya tiene deuda, y la mayoría se endeuda para el día a día. Endeudarse bien se aprende antes.', VERDE),
         ('EL JUEGO', 'Vivimos un año de decisiones: cuotas, mínimos, imprevistos, aguinaldo. El modelo PAMY y la cuenta cuota × cuotas.', CELESTE),
         ('LA HUELLA', 'BCU y Clearing: los dos espejos. Cada clavo marca 5 años; el banco los mira a los dos.', VIOLETA),
         ('DECIDIR BIEN', 'Con la IA reflexionamos cómo nos preparamos para un préstamo. La introspección también es una herramienta financiera.', NAVY)]
for i, (t, d, col) in enumerate(ideas):
    x = 0.6 + (i % 2) * 6.25; y = 1.85 + (i // 2) * 2.25
    b = rbox(s, x, y, 5.95, 2.0, CREMA, line=col, lw=2.5)
    text(s, x + 0.25, y + 0.18, 5.45, 1.7, [
        {'t': t, 's': 14, 'f': DISP, 'c': col, 'sa': 5},
        {'t': d, 's': 12.8, 'b': True},
    ])
text(s, 0.6, 6.45, 12.1, 0.6, [{'t': 'Cada participante elige UNA idea que no quiere olvidar… y se la dice a alguien. Decirlo es creerlo.', 's': 14.5, 'b': True, 'c': ROJO}])
footer(s)

s = slide(); head(s, 'TICKET DE SALIDA', 'Se llevan el juego (y su huella)')
b = rbox(s, 0.6, 1.85, 5.95, 4.3, CREMA, line=VERDE, lw=2)
if os.path.exists(QR1):
    s.shapes.add_picture(QR1, Inches(1.0), Inches(2.2), Inches(2.1), Inches(2.1))
text(s, 3.3, 2.4, 3.1, 2.2, [
    {'t': 'LLEVATE EL JUEGO', 's': 15, 'f': DISP, 'c': VERDE, 'sa': 5},
    {'t': 'Escanealo y jugalo cuando quieras, con quien quieras. Funciona sin internet.', 's': 12.5, 'b': True},
])
text(s, 1.0, 4.55, 5.2, 0.5, [{'t': 'chelabsweb.github.io/endeudarse-bien-stand', 's': 12, 'b': True, 'c': VERDE}])
b = rbox(s, 6.8, 1.85, 5.95, 4.3, CREMA, line=CELESTE, lw=2)
if os.path.exists(QR2):
    s.shapes.add_picture(QR2, Inches(7.2), Inches(2.2), Inches(2.1), Inches(2.1))
text(s, 9.5, 2.4, 3.0, 2.2, [
    {'t': 'TU HUELLA REAL', 's': 15, 'f': DISP, 'c': CELESTE, 'sa': 5},
    {'t': 'Consultá tu huella de verdad, gratis, al llegar a casa.', 's': 12.5, 'b': True},
])
text(s, 7.2, 4.55, 5.2, 0.5, [{'t': 'consultadeuda.bcu.gub.uy', 's': 12, 'b': True, 'c': CELESTE}])
footer(s)

s = slide(NAVY); dots(s, 11.6, 0.6, VERDE); dots(s, 0.8, 5.2, VIOLETA)
text(s, 1.0, 2.3, 11.3, 2.8, [
    {'t': 'LA IDEA FUERZA', 's': 18, 'f': DISP, 'c': VERDE, 'sa': 10},
    {'t': 'Un préstamo es un compromiso de ahorro,\nno una solución inmediata.', 's': 32, 'f': DISP, 'c': CREMA, 'sa': 10},
    {'t': 'Reflexionar, planificar y actuar con conciencia: las claves de un endeudamiento saludable y una huella financiera positiva.', 's': 15, 'b': True, 'c': BLANCO},
])
pin(s, 11.2, 6.3)
footer(s, dark=True)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Formacion_de_Formadores_Endeudarse_Bien_2026.pptx')
prs.save(out)
print('OK →', out, f'({NPAG[0]} slides)')
